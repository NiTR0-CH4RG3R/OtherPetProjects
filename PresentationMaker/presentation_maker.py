import sys

filename = ''

if len(sys.argv) == 2:
    filename = sys.argv[1]
else:
    sys.stderr.write('Usage : {0} <input file>\n'.format(sys.argv[0]))
    exit()

#Importing External libraries
from xlrd import open_workbook #This is the library for working with spreadsheets
from pptx import Presentation # All the Python-PPTX libraries are for deal with presentations
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


class Student:
    # subject is a dictionary
    def __init__(self,ind_no,name,subjects,total,avg,rank):
        self.ind_no = ind_no
        self.name = name
        self.subjects = subjects
        self.total = total
        self.avg = avg
        self.rank = rank

#Importing spreadsheet
wb = open_workbook(filename) 


default_pos = (0, 0.5) #Defualt position of a chart
default_size = (5,3.25) #Defualt size of a chart

#Function that generate a student from given data
def create_student(sheet,r,total_col_num,avg_col_num,rank_col_num,primary_subjects):
    index_no = sheet.cell(r,0).value
    name = sheet.cell(r,1).value
    subjects = {}
    total = sheet.cell(r,total_col_num).value
    avg = sheet.cell(r,avg_col_num).value
    rank = sheet.cell(r,rank_col_num).value
    for key in primary_subjects.keys():
        s = sheet.cell(r,primary_subjects[key]).value
        try:
            subjects.update({key:int(s)})
        except:
            if s == 'AB':
                subjects.update({key:0})
            if s == '' or s == None:
                subjects.update({key:0})
            else:
                pass  
    return Student(int(index_no),name,subjects,total,avg,int(rank))

#This function genrate chart data we need
def generate_chart_data(student):
        chart_data = CategoryChartData()
        chart_data.categories = list(student.subjects.keys())
        chart_data.add_series('Series 1',tuple(student.subjects.values()))
        return chart_data

for sheetname in wb.sheet_names():

    sheet = wb.sheet_by_name(sheetname)

    # Primary_Subjects = {Subject Name : Colum Number}
    # Subject Name should be a string
    # Column Number should be an intiger
    primary_subjects = {}
    students = []
    total_col_num = 0
    avg_col_num = 0
    rank_col_num = 0

    #Add all the subjects to primary_subject dictionary
    for c in range(2,sheet.ncols):
        if sheet.cell(0,c).value == 'Total':
            total_col_num, avg_col_num, rank_col_num = c, c+1, c+2
            break
        primary_subjects.update({sheet.cell(0,c).value:c})
    
    for r in range(1,sheet.nrows):
        try:        
            student = create_student(sheet,r,total_col_num,avg_col_num,rank_col_num,primary_subjects)
            students.append(student)
        except:
            print ('[+]',sys.exc_info())
            continue

    #Creating the presentation 
    prs = Presentation()
    #Adding a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    k = 0
    for student in students:
        chart_data = generate_chart_data(student)

        # x - X position (in Inches) of the Current chart
        # y - Y position (in Inches) of the Current chart
        # cx - Lenght of the Current chart
        # cy - Width of the Current chart
        x, y, cx, cy = Inches(default_pos[0]), Inches(default_pos[1]), Inches(default_size[0]), Inches(default_size[1])

        if k == 0:
            #Adding a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            k += 1
        elif k == 1:
            x = Inches(default_pos[0] + default_size[0])
            k += 1
        elif k == 2:
            y = Inches(default_pos[1] + default_size[1])
            k +=1
        else:
            x = Inches(default_pos[0] + default_size[0])
            y = Inches(default_pos[1] + default_size[1])
            k = 0
        try:
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        except:
            print ('[+]',sys.exc_info())
            continue

    #Saving current presentation
    prs.save(sheetname + '.pptx')

