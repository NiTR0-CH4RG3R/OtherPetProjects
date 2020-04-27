import sys
from xlrd import open_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

class student:
    def __init__(self,name,chemistry,physics,bio,english):
        self.name = name
        self.chemistry = chemistry
        self.physics = physics
        self.bio = bio
        self.english = english

filename = ""
out_filename = ""

if len(sys.argv) == 3:
    filename = sys.argv[1]
    out_filename = sys.argv[2]
else:
    sys.stderr.write("Usage: {0} <input excell file> <output presentation file>\n".format(sys.argv[0]))
    exit()

prs = Presentation()


wb = open_workbook(filename)

sheet1 = wb.sheets()[0]


chartsizex = 5
chartsizey = 3.25
chartdefualtposx = 0
chartdefualtposy = 0.5

for r in range(1,sheet1.nrows,4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    students = []
    for r1 in range(4):
        if r+r1 >= sheet1.nrows:
            break
        vals = []
        for c in range(1,6):
            vals.append(sheet1.cell(r+r1,c).value)
        
        for v in range(1,5):
            try:
                vals[v] = int(vals[v])
            except:
                vals[v] = 0

        stdud = student(vals[0],vals[1],vals[2],vals[3],vals[4])
        students.append(stdud)
    a = 0
    for st in students:
        chart_data = CategoryChartData()
        chart_data.categories = ['Chemistry','Physics','Bio','English']

        chart_data.add_series('Series 1',(st.chemistry,st.physics,st.bio,st.english))

        x, y, cx, cy = Inches(chartdefualtposx), Inches(chartdefualtposy), Inches(chartsizex), Inches(chartsizey)
        
        if a == 0:
            x=x
            y=y
        elif a == 1:
            x = Inches(chartdefualtposx + chartsizex)
        elif a == 2:
            y = Inches(chartdefualtposy + chartsizey)
        elif a == 3:
            x = Inches(chartdefualtposx + chartsizex)
            y = Inches(chartdefualtposy + chartsizey)
        else:
            pass
        a += 1
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

prs.save(out_filename + '.pptx')

